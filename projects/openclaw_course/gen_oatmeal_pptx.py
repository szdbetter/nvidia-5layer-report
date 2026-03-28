#!/usr/bin/env python3
"""Generate Oatmeal-style HTML slides and PPTX for OpenClaw course."""

import os
from pathlib import Path

SLIDES_DIR = Path("/tmp/slides_oatmeal")
IMG_DIR = SLIDES_DIR / "images"
OUT_DIR = Path("/root/.openclaw/workspace/projects/openclaw_course")
PREVIEW_DIR = OUT_DIR
PPTX_OUT = OUT_DIR / "openclaw_course_oatmeal.pptx"

SLIDES = [
    {
        "n": "01",
        "title1": "你的第一个硅基合伙人",
        "title2": "我们今天不讲 AI，我们讲"分身"。",
        "bullets": ["这不是工具升级，这是物种进化", "从"你用它"到"它替你"", "从"聊天机器人"到"数字员工""]
    },
    {
        "n": "02",
        "title1": "你现在的 AI，只是个高级打字机",
        "title2": "ChatGPT 再聪明，也得你亲手干活。",
        "bullets": ["你查资料、你写 prompt、你复制粘贴", "它只动嘴，你跑断腿", "本质是"搜索引擎 + 打字机""]
    },
    {
        "n": "03",
        "title1": "老板的工作，是下发任务",
        "title2": "真正的管理者，只负责定义"要什么"",
        "bullets": ["你喝茶，员工跑腿", "你决策，团队执行", "这才是正常的雇佣关系"]
    },
    {
        "n": "04",
        "title1": "OpenClaw：LLM 是大脑，Agent 是手脚",
        "title2": "它不是聊天机器人，它是你的操作系统。",
        "bullets": ["大脑（LLM）负责思考拆解", "手脚（Agent）负责执行闭环", "你只说一句话，它全自动完成"]
    },
    {
        "n": "05",
        "title1": "它自己思考、拆解、派活、推进到完成",
        "title2": "你睡觉的时候，它在干活。",
        "bullets": ["自动唤醒子代理团队", "自动搜索、自动重试、自动纠错", "第二天早上，结果已在微信里"]
    },
    {
        "n": "06",
        "title1": "别人卖"智力"，我们卖"执行力"",
        "title2": "有脑无手是残废，有脑有手才是员工。",
        "bullets": ["普通 AI：出主意，你执行（需人盯）", "OpenClaw：出主意 + 干活（全自动）", "从"参谋"到"将军"的质变"]
    },
    {
        "n": "07",
        "title1": "它能做的，远超你想象",
        "title2": "信息收割、代码生成、内容生产、数据监控。",
        "bullets": ["24 小时监控市场/竞品/舆情", "自动写代码、部署、测试", "自动生成研报/文章/视频脚本"]
    },
    {
        "n": "08",
        "title1": "它不是神，它有自己的边界",
        "title2": "它需要清晰的目标，不能替你承担法律责任。",
        "bullets": ["不能替你承担法律/财务责任", "不能处理模糊、矛盾的目标", "需要"人"定义战略方向"]
    },
    {
        "n": "09",
        "title1": "三种人，现在就需要配置数字员工",
        "title2": "老板、创业者、超级个体。",
        "bullets": ["企业高管：决策杠杆最大化", "创业者/一人公司：1 人抵 10 人", "内容创作者：规模化生产 + 变现"]
    },
    {
        "n": "10",
        "title1": "2026，为什么是元年？",
        "title2": "政策首次写入政府工作报告，70% 普及率目标。",
        "bullets": ["2026 政府工作报告首次写入"智能体"", "腾讯/阿里/字节/华为全部下场", "全球市场规模 3 年翻 3 倍"]
    },
    {
        "n": "11",
        "title1": "真实案例：他们已经在赚钱",
        "title2": "一人公司 6 个月被 8000 万美金收购。",
        "bullets": ["Base44：1 人+AI，被 Wix 8000 万美金收购", "Arcads AI：6 人团队，人均创收 279 万美金", "国内：深圳 OPC 硬件创业，7 人团队千万融资"]
    },
    {
        "n": "12",
        "title1": "要么拥有分身，要么成为别人的工具",
        "title2": "2026 年，配置你的第一个数字员工。",
        "bullets": ["今天不布局，明年没机会", "选择平台比自建更重要", "扫码，开启你的数字员工之旅"]
    },
]

HTML_TEMPLATE = """<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
* {{ box-sizing: border-box; }}
body {{ width:1280px; height:720px; margin:0; padding:0; overflow:hidden; position:relative; background-color:#F8F8F8; font-family:'PingFang SC','Microsoft YaHei',sans-serif; }}
p,h1,h2,h3,h4 {{ margin:0; padding:0; }}
</style>
</head>
<body>
  <img src="{img_path}" style="position:absolute;left:0;top:0;width:680px;height:720px;object-fit:cover;">
  <div style="position:absolute;right:0;top:0;width:600px;height:720px;background-color:#F8F8F8;"></div>
  <div style="position:absolute;right:0;top:0;width:600px;height:10px;background-color:#FF6B35;"></div>
  <div style="position:absolute;right:20px;top:30px;width:560px;">
    <h1 style="font-size:36px;font-weight:900;color:#333333;line-height:1.35;">{title1}</h1>
    <h2 style="font-size:24px;font-weight:700;color:#FF6B35;line-height:1.35;margin-top:10px;">{title2}</h2>
  </div>
  <div style="position:absolute;right:20px;top:255px;width:560px;height:4px;background-color:#333333;"></div>
  <div style="position:absolute;right:20px;top:275px;width:560px;">
    <p style="font-size:20px;font-weight:700;color:#333333;line-height:1.7;margin-bottom:12px;">● {b1}</p>
    <p style="font-size:20px;font-weight:700;color:#333333;line-height:1.7;margin-bottom:12px;">● {b2}</p>
    <p style="font-size:20px;font-weight:700;color:#333333;line-height:1.7;">● {b3}</p>
  </div>
  <div style="position:absolute;right:20px;bottom:20px;"><p style="font-size:40px;font-weight:900;color:#FF6B35;opacity:0.25;">{n}</p></div>
  <div style="position:absolute;right:0;bottom:0;width:600px;height:6px;background-color:#FF6B35;"></div>
</body>
</html>"""

# Generate HTML files
for s in SLIDES:
    n = s["n"]
    img_path = str(IMG_DIR / f"img-{n}.jpg")
    html = HTML_TEMPLATE.format(
        img_path=img_path,
        title1=s["title1"],
        title2=s["title2"],
        b1=s["bullets"][0],
        b2=s["bullets"][1],
        b3=s["bullets"][2],
        n=n
    )
    html_path = SLIDES_DIR / f"slide-{n}.html"
    html_path.write_text(html, encoding="utf-8")
    print(f"HTML {n} done")

print("HTML generation complete")

# Generate PPTX
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io

    prs = Presentation()
    prs.slide_width = Emu(9144000)   # 1280px at 96dpi = 9144000 EMU
    prs.slide_height = Emu(5143500)  # 720px

    # EMU per pixel at 96dpi: 1px = 9525 EMU
    PX = 9525

    blank_layout = prs.slide_layouts[6]  # blank

    for s in SLIDES:
        n = s["n"]
        slide = prs.slides.add_slide(blank_layout)
        img_file = str(IMG_DIR / f"img-{n}.jpg")

        # Left image: 680px wide, full height
        slide.shapes.add_picture(img_file, 0, 0, 680*PX, 720*PX)

        # Right panel background (#F8F8F8) - add a rectangle
        from pptx.util import Pt
        bg = slide.shapes.add_shape(1, 680*PX, 0, 600*PX, 720*PX)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        bg.line.fill.background()

        # Top orange bar
        bar_top = slide.shapes.add_shape(1, 680*PX, 0, 600*PX, 10*PX)
        bar_top.fill.solid()
        bar_top.fill.fore_color.rgb = RGBColor(0xFF, 0x6B, 0x35)
        bar_top.line.fill.background()

        # Bottom orange bar
        bar_bot = slide.shapes.add_shape(1, 680*PX, 710*PX, 600*PX, 10*PX)
        bar_bot.fill.solid()
        bar_bot.fill.fore_color.rgb = RGBColor(0xFF, 0x6B, 0x35)
        bar_bot.line.fill.background()

        # Title1
        tf1 = slide.shapes.add_textbox(700*PX, 30*PX, 560*PX, 120*PX)
        tf1.text_frame.word_wrap = True
        p1 = tf1.text_frame.paragraphs[0]
        p1.text = s["title1"]
        p1.font.bold = True
        p1.font.size = Pt(28)
        p1.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Title2
        tf2 = slide.shapes.add_textbox(700*PX, 160*PX, 560*PX, 80*PX)
        tf2.text_frame.word_wrap = True
        p2 = tf2.text_frame.paragraphs[0]
        p2.text = s["title2"]
        p2.font.bold = True
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xFF, 0x6B, 0x35)

        # Divider line
        div = slide.shapes.add_shape(1, 700*PX, 255*PX, 560*PX, 4*PX)
        div.fill.solid()
        div.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
        div.line.fill.background()

        # Bullets
        tf3 = slide.shapes.add_textbox(700*PX, 270*PX, 560*PX, 340*PX)
        tf3.text_frame.word_wrap = True
        for bi, bullet in enumerate(s["bullets"]):
            if bi == 0:
                p = tf3.text_frame.paragraphs[0]
            else:
                p = tf3.text_frame.add_paragraph()
            p.text = f"● {bullet}"
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(10)

        # Page number watermark
        tf4 = slide.shapes.add_textbox(1220*PX, 660*PX, 80*PX, 50*PX)
        p4 = tf4.text_frame.paragraphs[0]
        p4.text = n
        p4.font.bold = True
        p4.font.size = Pt(28)
        p4.font.color.rgb = RGBColor(0xFF, 0x6B, 0x35)

        print(f"PPTX slide {n} done")

    prs.save(str(PPTX_OUT))
    print(f"PPTX saved: {PPTX_OUT}")

except ImportError as e:
    print(f"python-pptx not available: {e}")
    print("Install with: pip install python-pptx")
