#!/usr/bin/env python3
"""BABA研报白板风格PPT - Volcengine Seedream 3.0"""
import os, json, time, urllib.request, urllib.error
from pathlib import Path

API_KEY = os.environ.get("VOLCENGINE_API_KEY", "")
API_URL = "https://ark.cn-beijing.volces.com/api/v3/images/generations"
MODEL = "doubao-seedream-3-0-t2i-250415"
OUT_DIR = Path("/root/.openclaw/workspace/ppt_outputs/baba_whiteboard")
OUT_DIR.mkdir(parents=True, exist_ok=True)

PAGES = [
    ("cover", "Realistic photograph of a large office whiteboard. Bold black thick marker calligraphy Chinese title '阿里巴巴BABA深度研报' centered. Blue marker subtitle '10x Alpha猎手 · 2026-03-17'. Red marker accent circle around '✅ Alpha'. Hand-drawn red bull arrow icon in corner. Clean whiteboard background with slight reflections."),
    ("p1_dashboard", "Realistic photograph of office whiteboard with thick marker handwriting. Title in black marker '投资决策驾驶舱'. Blue marker data table showing: 当前价格 $136.85, 目标价 $205, 胜率 68%, 赔率 1:2.8, 止损 $115. Red circle around '合理分批建仓'. Hand-drawn chart arrow going up. Clean white background."),
    ("p2_scorecard", "Realistic photograph of whiteboard with marker writing. Black title '14维度评分卡'. Blue marker showing scores: 技术护城河 7.5, 商业模式 7.5, 收入增长 6.5, 估值合理性 8.5, 产品迭代 8.0. Red marker total '总分 7.16/10'. Hand-drawn radar chart sketch in blue and red markers. Clean whiteboard."),
    ("p3_logic", "Realistic photograph of whiteboard. Black marker title '核心投资逻辑'. Blue marker: 认知差分析. Two columns: left '市场认为' with crossed out text '增长停滞电商公司', right '实际上' with '中国最大AI基础设施'. Red underline under 'P/S=2.31 历史10%分位'. Hand-drawn lightbulb icon. Whiteboard background."),
    ("p4_segments", "Realistic photograph of whiteboard with hand-drawn pie chart in blue and red markers. Black title '业务板块'. Segments labeled: 淘宝天猫50%, 云智能12%+34%增速, 国际商业11%+33%, 菜鸟10%, 其他17%. Red star next to 云智能. Blue annotations. Clean whiteboard."),
    ("p5_ai", "Realistic photograph of whiteboard. Black marker title 'AI战略 Qwen生态'. Blue marker data: 200+模型, 3亿下载, 397B MoE, Apache 2.0. Red marker '超GPT-5.2'. Hand-drawn neural network diagram in blue marker. Red circle around '$52B投资'. Whiteboard background."),
    ("p6_competition", "Realistic photograph of whiteboard with two hand-drawn comparison tables in markers. Left table '电商竞争' black title: 阿里44%, PDD 19%, JD 24%, 抖音⚠️. Right table '云计算' blue title: 阿里36%#1, 华为15%, 腾讯12%. Red warning triangle next to 抖音. Clean whiteboard."),
    ("p7_risks", "Realistic photograph of whiteboard. Red marker title '风险与熔断'. Black marker list: 估值熔断<$115, 质量熔断EBITA<11%, 增长熔断云<15%, 人才熔断Qwen离职, 监管黑天鹅. Red circles around critical items. Hand-drawn warning sign. Whiteboard background."),
    ("p8_strategy", "Realistic photograph of whiteboard. Black marker title '操作策略'. Blue marker: 建仓$130-138, 加仓$138-155, 目标$192-205. Red marker: 止损$115, 仓位5-8%. Hand-drawn upward staircase chart. Red star next to '3-19财报'. Green checkmark. Clean whiteboard."),
]

def generate_image(name, prompt):
    out_path = OUT_DIR / f"{name}.jpg"
    if out_path.exists():
        print(f"[skip] {name} exists")
        return str(out_path)
    
    data = json.dumps({
        "model": MODEL,
        "prompt": prompt,
        "n": 1,
        "size": "1792x1024",
        "response_format": "url"
    }).encode("utf-8")
    
    req = urllib.request.Request(API_URL, data=data, headers={
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }, method="POST")
    
    print(f"[gen] {name}...", end=" ", flush=True)
    try:
        resp = urllib.request.urlopen(req, timeout=120)
        result = json.loads(resp.read())
        url = result["data"][0]["url"]
        # download
        urllib.request.urlretrieve(url, str(out_path))
        print(f"✅ saved")
        return str(out_path)
    except Exception as e:
        print(f"❌ {e}")
        return None

def build_pptx(image_paths):
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for img in image_paths:
        if img and Path(img).exists():
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    
    out = str(OUT_DIR / "BABA_Whiteboard_2026-03-17.pptx")
    prs.save(out)
    print(f"\n[pptx] {out}")
    return out

if __name__ == "__main__":
    images = []
    for name, prompt in PAGES:
        img = generate_image(name, prompt)
        images.append(img)
        time.sleep(1)  # rate limit
    
    pptx_path = build_pptx(images)
    print(f"DONE: {pptx_path}")
