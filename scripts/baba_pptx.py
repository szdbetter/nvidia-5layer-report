#!/usr/bin/env python3
"""BABA研报 → PPT生成器"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
BLUE = RGBColor(0x1E, 0x3A, 0x5F)
RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF0)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
GREEN = RGBColor(0x27, 0xAE, 0x60)
GOLD = RGBColor(0xF3, 0x9C, 0x12)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = "Noto Sans CJK SC"
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_table(slide, left, top, rows_data, col_widths, header_color=BLUE):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          sum(Inches(w) for w in col_widths), Inches(rows * 0.45))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Noto Sans CJK SC"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_color
                else:
                    p.font.color.rgb = BLACK
    return table_shape

# ===== 封面 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK_BG)
add_text(s, 1, 1.5, 11, 1.5, "阿里巴巴（BABA）深度研报", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(s, 1, 3.2, 11, 0.8, "10x Alpha 猎手 · 第二阶段 · v5.0", 24, GOLD, False, PP_ALIGN.CENTER)
add_text(s, 1, 4.5, 11, 0.6, "OpenClaw AI Research Engine | 2026-03-17", 16, RGBColor(0xAA,0xAA,0xAA), False, PP_ALIGN.CENTER)
add_text(s, 1, 5.5, 11, 0.6, "加权总分 7.16/10 → ✅ Alpha | 胜率 68% | 赔率 1:2.8", 20, GREEN, True, PP_ALIGN.CENTER)

# ===== P1: 投资决策驾驶舱 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_text(s, 0.5, 0.3, 12, 0.8, "🏁 投资决策驾驶舱", 32, BLUE, True)
data = [
    ["项目", "结论"],
    ["投资建议", "合理分批建仓（Moderate Accumulate）"],
    ["时机判断", "错杀修复期 — 反垄断清零+AI战略落地+估值历史低位"],
    ["触发事件", "2026-03-19 Q3 FY2026财报"],
    ["当前价格", "$136.85（极度低估区间）"],
    ["胜率预测", "68%"],
    ["赔率预测", "1:2.8（目标$192-205，止损$118）"],
]
add_table(s, 0.8, 1.5, data, [3, 9])
add_text(s, 0.8, 5.5, 11, 0.8, "核心逻辑：中国云AI绝对领导者 + 电商低估值护城河 + 反垄断出清 → 双击估值+盈利修复", 16, RED, True)

# ===== P2: 14维度评分卡 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_text(s, 0.5, 0.3, 12, 0.8, "📊 14维度评分卡（加权总分 7.16/10）", 32, BLUE, True)
scores = [
    ["维度", "分数", "权重得分"],
    ["1.技术护城河(12%)", "7.5", "0.90"],
    ["2.商业模式(10%)", "7.5", "0.75"],
    ["3.收入增长(10%)", "6.5", "0.65"],
    ["4.盈利能力(8%)", "7.5", "0.60"],
    ["5.现金流(8%)", "7.5", "0.60"],
    ["6.市场地位(8%)", "7.0", "0.56"],
    ["7.管理团队(6%)", "6.0", "0.36"],
    ["8.股东结构(6%)", "6.0", "0.36"],
    ["9.估值合理性(10%)", "8.5", "0.85"],
    ["10.政策监管(6%)", "6.0", "0.36"],
    ["11.产品迭代(6%)", "8.0", "0.48"],
    ["12.客户基础(4%)", "7.5", "0.30"],
    ["13.供应链(4%)", "5.5", "0.22"],
    ["14.分析师共识(2%)", "8.5", "0.17"],
]
add_table(s, 0.8, 1.3, scores, [4.5, 2, 2.5])

# ===== P3: 核心投资逻辑 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_text(s, 0.5, 0.3, 12, 0.8, "💡 核心投资逻辑", 32, BLUE, True)
add_text(s, 0.8, 1.5, 11, 1, "认知差：市场定价阿里是增长停滞的电商公司", 22, BLACK, True)
add_text(s, 0.8, 2.3, 11, 1.5, "• P/S = 2.31x，较10年中位数6.74低66%，处历史10-15%分位\n• 市场忽略了云AI业务的爆发性增长（34%增速，AI三位数）\n• 反垄断整改2024年8月正式完成，监管风险已清零", 16, BLACK)
add_text(s, 0.8, 4.0, 11, 1, "实际上它是：中国最大AI基础设施运营商", 22, RED, True)
add_text(s, 0.8, 4.8, 11, 1.5, "• Qwen 3.5 超 GPT-5.2（Math-Vision基准），全球3亿+下载\n• 电商 = 现金牛（$27B/年EBITA） → 持续输血云AI飞轮\n• 对标 Microsoft 2016-2020 云转型路径，估值扩张空间巨大", 16, BLACK)

# ===== P4: 业务板块拆解 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_text(s, 0.5, 0.3, 12, 0.8, "📦 业务板块拆解", 32, BLUE, True)
biz = [
    ["业务", "营收(估算)", "占比", "增速", "角色"],
    ["淘宝天猫(TTG)", "~500B元", "~50%", "+低单位数", "现金牛🐄"],
    ["云智能集团", "~115B元", "~12%", "+34%", "增长引擎🚀"],
    ["国际数字商业", "~108B元", "~11%", "+33%", "全球化扩张"],
    ["菜鸟物流", "~100B元", "~10%", "+中单位数", "战略支撑"],
    ["本地生活/大文娱", "~170B元", "~17%", "混合", "生态补充"],
    ["合计", "~996B元", "100%", "+5-6%", "—"],
]
add_table(s, 0.5, 1.3, biz, [3, 2, 1.5, 2, 3])

# ===== P5: AI战略 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_text(s, 0.5, 0.3, 12, 0.8, "🤖 AI战略：Qwen生态", 32, BLUE, True)
ai = [
    ["指标", "数据"],
    ["Qwen系列模型", "200+ 开源模型"],
    ["全球下载量", "3亿+"],
    ["衍生模型", "10万+"],
    ["Qwen 3.5 参数", "397B MoE架构"],
    ["许可证", "Apache 2.0（商用自由）"],
    ["上下文窗口", "256K tokens"],
    ["vs GPT-5.2", "超越（Math-Vision基准）"],
    ["成本优化", "较上代降60%，性能提升8x"],
    ["3年AI投资", "3800亿元($52B)"],
]
add_table(s, 0.8, 1.3, ai, [4, 7])

# ===== P6: 竞争格局 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_text(s, 0.5, 0.3, 12, 0.8, "⚔️ 竞争格局矩阵", 32, BLUE, True)
add_text(s, 0.5, 1.2, 5.5, 0.5, "中国电商竞争", 20, RED, True)
ecom = [
    ["平台", "份额", "威胁度"],
    ["淘宝天猫", "~44%", "—"],
    ["拼多多", "~19%", "中"],
    ["JD.com", "~24%", "低"],
    ["抖音电商", "快速增长", "⚠️高"],
]
add_table(s, 0.5, 1.8, ecom, [2.5, 1.5, 1.5])
add_text(s, 6.5, 1.2, 5.5, 0.5, "中国云计算竞争", 20, RED, True)
cloud = [
    ["云服务商", "份额", "AI能力"],
    ["阿里云", "36% #1", "Qwen系列"],
    ["华为云", "~15%", "昇腾芯片"],
    ["腾讯云", "~12%", "混元大模型"],
    ["字节云", "快速增长", "豆包"],
]
add_table(s, 6.5, 1.8, cloud, [2.5, 1.5, 2.5])

# ===== P7: 风险与熔断 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_text(s, 0.5, 0.3, 12, 0.8, "⚠️ 风险矩阵与熔断条件", 32, RED, True)
risks = [
    ["熔断类型", "触发条件", "风险分"],
    ["估值熔断", "股价<$115 或 P/S<1.8x", "🔴高"],
    ["质量熔断", "EBITA利润率连续两季<11%", "🟡中"],
    ["增长熔断", "云收入增速连续两季<15%", "🟡中"],
    ["人才熔断", "Qwen核心负责人离职", "🟡中"],
    ["监管黑天鹅", "BABA从NYSE摘牌", "🔴尾部"],
    ["竞争熔断", "云份额跌破25%（连续两季）", "🟡中"],
]
add_table(s, 0.5, 1.3, risks, [3, 5.5, 2], RED)

# ===== P8: 操作策略 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_text(s, 0.5, 0.3, 12, 0.8, "🎯 操作策略与监控", 32, BLUE, True)
ops = [
    ["操作", "价格", "说明"],
    ["首次建仓", "$130-138", "当前区间，分批2-3次"],
    ["加仓(财报后)", "$138-155", "若云收入>预期"],
    ["目标止盈(12月)", "$192-205", "分析师共识"],
    ["激进目标(18月)", "$240-260", "P/S修复至4-4.5x"],
    ["止损线", "$115", "跌破则逻辑受损"],
    ["仓位建议", "5-8%", "中国科技主力配置"],
]
add_table(s, 0.5, 1.3, ops, [3, 2.5, 5.5])
add_text(s, 0.5, 5.0, 12, 0.5, "🔥 关键催化剂：2026-03-19 Q3 FY2026财报（3天后）", 20, RED, True)
add_text(s, 0.5, 5.6, 12, 0.8, "重点关注：云收入增速(目标>20%) | AI收入占比 | TTG EBITA利润率 | FCF变化 | capex指引", 14, BLACK)

# ===== 尾页 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK_BG)
add_text(s, 1, 2, 11, 1.5, "BABA · 7.16/10 · ✅ Alpha", 44, GREEN, True, PP_ALIGN.CENTER)
add_text(s, 1, 3.8, 11, 1, "胜率68% | 赔率1:2.8 | 目标$205 | 止损$115", 24, GOLD, False, PP_ALIGN.CENTER)
add_text(s, 1, 5.2, 11, 0.8, "⚠️ 免责声明：本报告仅供研究参考，不构成投资建议", 14, RGBColor(0x99,0x99,0x99), False, PP_ALIGN.CENTER)
add_text(s, 1, 6.0, 11, 0.5, "OpenClaw AI Research Engine | 2026-03-17", 12, RGBColor(0x77,0x77,0x77), False, PP_ALIGN.CENTER)

# 保存
out = "/root/.openclaw/workspace/ppt_outputs/baba_report/BABA_Research_2026-03-17.pptx"
prs.save(out)
print(f"SAVED: {out}")
